from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import pandas as pd
import os

# Import your modules
from data_loader import load_excel_data
from aggregator import generate_summaries
from visualizer import (
    plot_daily_summary,
    plot_weekly_summary,
    plot_common_issues,
    plot_engineer_workload
)
from genai_insights import generate_summary_text

# Load data
file_path = '../sample_data/demo_file.xlsx'
df = load_excel_data(file_path)
df['Date'] = pd.to_datetime(df['Date'])

# Generate summaries
daily_summary, weekly_summary, monthly_summary, common_issues, engineer_workload = generate_summaries(df)

# Save plots
plot_daily_summary(daily_summary)
plt.savefig("daily_summary.png")
plt.clf()

plot_weekly_summary(weekly_summary)
plt.savefig("weekly_summary.png")
plt.clf()

plot_common_issues(common_issues)
plt.savefig("common_issues.png")
plt.clf()

plot_engineer_workload(engineer_workload)
plt.savefig("engineer_workload.png")
plt.clf()

# Generate summary text
summary_text = generate_summary_text(daily_summary, weekly_summary, common_issues, engineer_workload)

# Create PowerPoint
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Client Issue Summary Report"
slide.placeholders[1].text = "Generated using GenSight"

# Summary Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary Insights"
slide.placeholders[1].text = summary_text

# Function to add image slides
def add_image_slide(title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

# Add image slides
add_image_slide("Daily Summary", "daily_summary.png")
add_image_slide("Weekly Summary", "weekly_summary.png")
add_image_slide("Common Issues", "common_issues.png")
add_image_slide("Engineer Workload", "engineer_workload.png")

# Save presentation
prs.save("client_presentation.pptx")
print("✅ Presentation saved as client_presentation.pptx")
